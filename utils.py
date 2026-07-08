"""
utils.py
---------
Fonctions de calcul pour le dashboard de réduction du downtime.
Toutes les fonctions prennent un DataFrame pandas "standardisé" avec au minimum
les colonnes : Date, Machine, Cause, Duree_min
(la colonne Responsable est optionnelle)

Ce module est volontairement séparé de app.py pour que la logique de calcul
soit testable indépendamment de l'interface Streamlit.
"""

import pandas as pd
import numpy as np

REQUIRED_COLUMNS = ["Date", "Machine", "Cause", "Duree_min"]

# Mots-clés utilisés pour deviner automatiquement la colonne "Famille" (catégorie
# de machine : Cutting Machine, Kit Seal, Outils, Press...) dans un export RadGrid
# ou tout autre export de pannes. Chaque nouvel export peut avoir des noms de
# colonnes légèrement différents, donc on garde une liste large de candidats.
FAMILY_COLUMN_CANDIDATES = [
    "sub description", "famille", "family", "categorie", "catégorie",
    "category", "plant group", "type de machine", "asset type",
]


def auto_detect_family_column(columns):
    """Devine la colonne 'Famille' (catégorie de machine) dans une liste de colonnes,
    en cherchant les mots-clés les plus probables. Retourne None si rien ne correspond."""
    for candidate in FAMILY_COLUMN_CANDIDATES:
        for col in columns:
            if candidate in str(col).lower():
                return col
    return None


def _detect_header_row(raw: pd.DataFrame, max_scan: int = 10) -> int:
    """
    Cherche, dans les premières lignes du fichier brut, celle qui ressemble le plus
    à une ligne d'en-tête (plusieurs cellules remplies, contenant des mots-clés
    typiques d'un export de pannes comme 'date', 'duration'/'durée', 'asset', etc.).
    """
    keywords = ["date", "duration", "durée", "duree", "asset", "machine", "cause", "description"]
    for i in range(min(max_scan, len(raw))):
        row = raw.iloc[i]
        non_empty = row.notna().sum()
        row_text = " ".join(str(v).lower() for v in row if pd.notna(v))
        if non_empty >= 3 and any(k in row_text for k in keywords):
            return i
    return 0


def load_data(file) -> pd.DataFrame:
    """
    Charge un fichier Excel (.xlsx/.xls) ou CSV et retourne un DataFrame brut.

    Robuste aux différents exports RadGrid : certains fichiers ont l'en-tête
    directement en ligne 1, d'autres ont une ou plusieurs lignes de titre/logo
    au-dessus. Cette fonction détecte automatiquement la vraie ligne d'en-tête
    au lieu de supposer qu'elle est toujours en première position.
    """
    name = getattr(file, "name", str(file))

    if name.lower().endswith(".csv"):
        raw = pd.read_csv(file, header=None)
    elif name.lower().endswith(".xls"):
        raw = pd.read_excel(file, header=None, engine="xlrd")
    else:
        raw = pd.read_excel(file, header=None)

    header_row_idx = _detect_header_row(raw)

    df = raw.iloc[header_row_idx + 1:].copy()
    df.columns = raw.iloc[header_row_idx]

    # Retire les colonnes sans nom (souvent des colonnes vides en bout de fichier)
    df = df.loc[:, df.columns.notna()]
    df.columns = [str(c).strip() for c in df.columns]

    # Dé-doublonne les noms de colonnes identiques (ex: deux colonnes "Trade"
    # dans un export RadGrid). Sans ça, choisir l'une d'elles dans le mapping
    # de l'interface renomme les DEUX colonnes en même temps (pandas renomme
    # par label, pas par position), ce qui corrompt silencieusement les
    # données. On garde le 1er nom tel quel, et on suffixe les suivants.
    seen = {}
    new_cols = []
    for c in df.columns:
        if c not in seen:
            seen[c] = 0
            new_cols.append(c)
        else:
            seen[c] += 1
            new_cols.append(f"{c} ({seen[c]+1})")
    df.columns = new_cols

    df = df.dropna(how="all").reset_index(drop=True)
    return df


def standardize_columns(df: pd.DataFrame, mapping: dict) -> pd.DataFrame:
    """
    Renomme les colonnes du fichier utilisateur vers le schéma standard
    attendu par les fonctions de calcul, à partir d'un mapping choisi
    dans l'interface (ex: {"Date panne": "Date", "Équipement": "Machine", ...}).
    """
    df = df.rename(columns=mapping)
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    if missing:
        raise ValueError(
            f"Colonnes manquantes après mapping : {missing}. "
            f"Merci de vérifier la correspondance des colonnes."
        )
    return df


def _fix_ambiguous_dates(dates: pd.Series, max_ecart_jours: int = 45) -> tuple:
    """
    Corrige automatiquement les dates dont le jour et le mois ont été inversés
    à l'export (bug très fréquent avec les exports RadGrid/Excel : une partie
    des lignes garde une date texte "correcte" pendant qu'une autre partie a
    été re-sérialisée par Excel avec la mauvaise convention jour/mois, par ex.
    "01/07/2026" (1 juillet, format JJ/MM) réinterprété comme le 7 janvier
    (format MM/JJ)).

    Principe (indépendant du fichier, ne suppose aucune plage de dates connue
    à l'avance) :
      1. On sépare les dates "sûres" (jour > 12 : aucune ambiguïté possible
         jour/mois) des dates "ambiguës" (jour <= 12 ET mois <= 12 : les deux
         lectures sont valides).
      2. On calcule la date médiane du groupe "sûr" comme référence de la
         période réelle des données.
      3. Toute date ambiguë située à plus de `max_ecart_jours` jours de cette
         référence est retestée en inversant jour et mois. Si l'inversion la
         rapproche nettement de la référence, on l'applique.
      4. S'il n'y a pas assez de dates "sûres" pour établir une référence
         fiable (ex: fichier couvrant un seul mois), on ne touche à rien :
         mieux vaut ne pas corriger que corriger au hasard.

    Retourne (dates_corrigees: pd.Series, nb_corrections: int).
    """
    dates = dates.copy()
    valides = dates.dropna()
    surs = valides[valides.dt.day > 12]

    if len(surs) < 3:
        return dates, 0

    reference = surs.median()
    n_corrected = 0

    for idx, d in valides.items():
        if d.day > 12:
            continue  # non ambigu, on ne touche pas
        ecart_original = abs((d - reference).days)
        if ecart_original <= max_ecart_jours:
            continue  # déjà cohérent avec le reste des données
        try:
            d_inverse = d.replace(month=d.day, day=d.month)
        except ValueError:
            continue  # inversion impossible (ex: mois=13), on laisse tel quel
        if abs((d_inverse - reference).days) < ecart_original:
            dates.at[idx] = d_inverse
            n_corrected += 1

    return dates, n_corrected


def clean_data(df: pd.DataFrame) -> pd.DataFrame:
    """Nettoie les types (dates, durées numériques) et retire les lignes invalides.

    Corrige aussi automatiquement les dates ambiguës (jour/mois inversés),
    un bug d'export très fréquent qui, sinon, disperse silencieusement une
    partie des pannes dans de fausses semaines (ex: des pannes de la semaine
    27 qui se retrouvent affichées en semaine 2, 6 ou 10)."""
    df = df.copy()
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    df["Duree_min"] = pd.to_numeric(df["Duree_min"], errors="coerce")
    df["Machine"] = df["Machine"].astype(str).str.strip()
    df["Cause"] = df["Cause"].astype(str).str.strip()

    # La colonne "Famille" (catégorie de machine) est optionnelle : on la nettoie
    # seulement si elle a été mappée, sans jamais l'exiger.
    if "Famille" in df.columns:
        df["Famille"] = df["Famille"].astype(str).str.strip()

    df["Date"], n_dates_corrigees = _fix_ambiguous_dates(df["Date"])

    n_before = len(df)
    df = df.dropna(subset=["Date", "Duree_min"])
    df = df[df["Duree_min"] >= 0]
    n_after = len(df)

    df.attrs["lignes_supprimees"] = n_before - n_after
    df.attrs["dates_corrigees"] = n_dates_corrigees
    return df


def compute_pareto(df: pd.DataFrame, group_col: str = "Cause") -> pd.DataFrame:
    """Pareto niveau 1 : classement des causes (ou machines) par durée totale cumulée."""
    grouped = df.groupby(group_col)["Duree_min"].sum().sort_values(ascending=False)
    pct = (grouped / grouped.sum()) * 100
    cumul = pct.cumsum()
    result = pd.DataFrame({
        "Duree_totale_min": grouped,
        "Pourcentage_%": pct.round(1),
        "Cumul_%": cumul.round(1),
    })
    result.index.name = group_col
    return result.reset_index()


def compute_pareto_level2(df: pd.DataFrame, top_value: str, level1_col: str = "Cause",
                           level2_col: str = "Machine") -> pd.DataFrame:
    """Pareto niveau 2 : détail par machine (ou autre dimension) pour UNE cause donnée."""
    subset = df[df[level1_col] == top_value]
    if subset.empty:
        return pd.DataFrame(columns=[level2_col, "Duree_totale_min", "Pourcentage_%", "Cumul_%"])
    return compute_pareto(subset, group_col=level2_col)


def compute_frequency_pareto(df: pd.DataFrame, group_col: str = "Machine") -> pd.DataFrame:
    """
    Pareto par FRÉQUENCE : classement par NOMBRE de pannes (et non par durée cumulée),
    comme le tableau « Frequent breakdowns » utilisé en interne chez Versigent.

    Complémentaire du Pareto par durée (compute_pareto) : un équipement peut apparaître
    en tête ici même si chaque panne est courte, dès lors qu'il tombe en panne très
    souvent — ce qui révèle un problème de fiabilité récurrent, différent d'un simple
    gros arrêt ponctuel repéré par le Pareto en durée.

    Retourne les mêmes noms de colonnes que compute_pareto (Pourcentage_%/Cumul_%),
    pour pouvoir réutiliser exactement les mêmes fonctions de graphique/export.
    """
    grouped = df.groupby(group_col).size().sort_values(ascending=False)
    grouped.name = "Nombre_pannes"
    pct = (grouped / grouped.sum()) * 100
    cumul = pct.cumsum()
    result = pd.DataFrame({
        "Nombre_pannes": grouped,
        "Pourcentage_%": pct.round(1),
        "Cumul_%": cumul.round(1),
    })
    result.index.name = group_col
    return result.reset_index()


def compute_paynter(df: pd.DataFrame, group_col: str = "Cause") -> pd.DataFrame:
    """Tableau Paynter : durée totale par cause et par semaine ISO."""
    df = df.copy()
    df["Semaine"] = df["Date"].dt.isocalendar().year.astype(str) + "-S" + \
        df["Date"].dt.isocalendar().week.astype(str).str.zfill(2)
    paynter = df.pivot_table(
        index=group_col, columns="Semaine", values="Duree_min",
        aggfunc="sum", fill_value=0
    )
    # Trie les semaines chronologiquement
    paynter = paynter[sorted(paynter.columns)]
    return paynter


def compute_alerts(df: pd.DataFrame, threshold_min: float, group_col: str = "Machine") -> pd.DataFrame:
    """Retourne les machines (ou causes) dont le cumul de downtime dépasse le seuil."""
    totals = df.groupby(group_col)["Duree_min"].sum().sort_values(ascending=False)
    alerts = totals[totals > threshold_min]
    return alerts.reset_index().rename(columns={"Duree_min": "Duree_totale_min"})


def predict_next_week(paynter: pd.DataFrame, row_name: str) -> dict:
    """
    Prédiction simple de la durée de panne de la semaine suivante pour une cause donnée,
    par régression linéaire sur l'historique des semaines. Si moins de 3 points de
    données, retombe sur une moyenne mobile simple.
    """
    if row_name not in paynter.index:
        return {"prediction": None, "methode": "aucune donnée"}

    y = paynter.loc[row_name].values.astype(float)
    n = len(y)

    if n < 3:
        pred = float(np.mean(y)) if n > 0 else 0.0
        return {"prediction": round(pred, 1), "methode": "moyenne simple (peu de données)"}

    X = np.arange(n).reshape(-1, 1)
    try:
        from sklearn.linear_model import LinearRegression
        model = LinearRegression().fit(X, y)
        pred = model.predict([[n]])[0]
        pred = max(0.0, float(pred))  # une durée ne peut pas être négative
        return {"prediction": round(pred, 1), "methode": "régression linéaire"}
    except ImportError:
        # Fallback si scikit-learn n'est pas installé : moyenne mobile sur les 4 dernières semaines
        pred = float(np.mean(y[-4:]))
        return {"prediction": round(pred, 1), "methode": "moyenne mobile (scikit-learn indisponible)"}


def generate_text_summary(df: pd.DataFrame, pareto1: pd.DataFrame, group_col: str = "Cause") -> str:
    """
    Génère un court résumé en français, directement réutilisable dans un rapport
    ou une présentation, à partir des résultats calculés.
    """
    if df.empty or pareto1.empty:
        return "Pas assez de données pour générer un résumé."

    nb_events = len(df)
    total_h = round(df["Duree_min"].sum() / 60, 1)
    top_row = pareto1.iloc[0]
    top_name = top_row[group_col]
    top_pct = top_row["Pourcentage_%"]
    top_2_cumul = pareto1.iloc[min(1, len(pareto1) - 1)]["Cumul_%"]

    nb_causes_80 = (pareto1["Cumul_%"] <= 80).sum() + 1
    nb_causes_80 = min(nb_causes_80, len(pareto1))

    resume = (
        f"Sur la période analysée, {nb_events} événements de panne ont été enregistrés, "
        f"représentant un total de {total_h} heures d'arrêt. "
        f"La cause principale est « {top_name} », responsable à elle seule de {top_pct}% "
        f"du downtime total. "
        f"Conformément à la règle de Pareto (80/20), {nb_causes_80} cause(s) sur {len(pareto1)} "
        f"expliquent déjà {top_2_cumul}% du problème. "
        f"Il est donc recommandé de concentrer les actions correctives en priorité sur « {top_name} » "
        f"avant de traiter les causes secondaires."
    )
    return resume


def build_action_plan_template(pareto1: pd.DataFrame, group_col: str = "Cause", top_n: int = 5) -> pd.DataFrame:
    """
    Génère un plan d'action pré-rempli avec les causes prioritaires (issues du Pareto),
    dans le même format que le tableau "Action Plan" utilisé en interne chez Versigent
    (Primary Reason Code / Secondary Reason Code + Root Cause / Corrective Action / Effect).
    Les colonnes à remplir manuellement par l'équipe maintenance sont laissées vides.
    """
    if pareto1.empty:
        return pd.DataFrame()

    top = pareto1.head(top_n).copy()
    plan = pd.DataFrame({
        "Priorité": range(1, len(top) + 1),
        "Primary Reason Code": top[group_col],
        "Root Cause Effect (%)": top["Pourcentage_%"],
        "Secondary Reason Code / Root Cause (à compléter)": "",
        "Permanent Corrective Action (à compléter)": "",
        "Responsable (à compléter)": "",
        "Date échéance (à compléter)": "",
        "Statut (à compléter)": "À traiter",
    })
    return plan


def build_action_plan_par_famille(pareto_par_famille_detail: dict, top_n: int = 3) -> pd.DataFrame:
    """
    Construit un plan d'action pré-rempli à partir des équipements les plus critiques
    de CHAQUE famille (top_n par famille : Cutting Machine, Kit Seal, Outils, Press...),
    dans le même esprit que la feuille "ACTION PLAN" utilisée en interne chez Versigent
    (qui liste des équipements précis comme M02, V827G, Crimping BT 752 ST01 — pas des
    codes de cause génériques).

    À utiliser à la place de build_action_plan_template quand une colonne Famille est
    disponible : c'est ce que l'entreprise attend (un plan d'action qui ne mélange pas
    toutes les familles, comme le Pareto lui-même).
    """
    rows = []
    item_nr = 1
    for famille, sous_pareto in pareto_par_famille_detail.items():
        if sous_pareto.empty:
            continue
        top = sous_pareto.head(top_n)
        machine_col = sous_pareto.columns[0]
        for _, r in top.iterrows():
            rows.append({
                "Item Nr": item_nr,
                "Famille": famille,
                "Équipement / Machine": r[machine_col],
                "Durée totale (min)": r["Duree_totale_min"],
                "% du downtime de la famille": r["Pourcentage_%"],
                "Cause racine (à compléter)": "",
                "Action corrective (à compléter)": "",
                "Responsable (à compléter)": "",
                "Date échéance (à compléter)": "",
                "Statut (à compléter)": "À traiter",
            })
            item_nr += 1
    return pd.DataFrame(rows)


def format_excel_sheet(worksheet, n_cols, header_color="1F4E78"):
    """
    Applique une mise en forme professionnelle à une feuille Excel :
    en-tête coloré et en gras, texte blanc, colonnes ajustées à la largeur du contenu,
    première ligne figée. Rend le fichier exporté directement présentable en réunion.
    """
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    header_fill = PatternFill(start_color=header_color, end_color=header_color, fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)

    for col_idx in range(1, n_cols + 1):
        cell = worksheet.cell(row=1, column=col_idx)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Ajuste automatiquement la largeur des colonnes selon leur contenu
    for col_idx in range(1, n_cols + 1):
        letter = get_column_letter(col_idx)
        max_len = 10
        for row in worksheet.iter_rows(min_col=col_idx, max_col=col_idx):
            for cell in row:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value)))
        worksheet.column_dimensions[letter].width = min(max_len + 3, 45)

    worksheet.freeze_panes = "A2"


def add_seuil80_column(worksheet, n_rows, table_col_count, start_row=1, seuil_value=80):
    """
    Ajoute une colonne "Seuil_80" (valeur constante 80) juste à droite d'un tableau de
    Pareto déjà écrit dans la feuille (colonnes 1..table_col_count, lignes
    start_row+1..start_row+n_rows pour les données, start_row pour l'en-tête).
    Sert à tracer la ligne de seuil 80% dans le graphique Pareto associé
    (add_pareto_excel_chart). Retourne l'index de colonne de la nouvelle colonne.
    """
    seuil_col = table_col_count + 1
    worksheet.cell(row=start_row, column=seuil_col, value="Seuil_80")
    for i in range(1, n_rows + 1):
        worksheet.cell(row=start_row + i, column=seuil_col, value=seuil_value)
    return seuil_col


def add_pareto_excel_chart(worksheet, header_row, n_rows, cat_col, val_col, cumul_col,
                            title, anchor, seuil_col=None, value_axis_title="Durée totale (min)",
                            width=17, height=9):
    """
    Insère un vrai graphique Pareto Excel natif (barres = valeur + courbe de cumul %
    sur axe secondaire, avec seuil 80% en pointillé si `seuil_col` est fourni) dans la
    feuille `worksheet`, à partir d'un tableau déjà écrit (en-tête à `header_row`,
    `n_rows` lignes de données juste en dessous). `cat_col`/`val_col`/`cumul_col`/
    `seuil_col` sont des index de colonnes 1-based (1=A, 2=B...), typiquement ceux
    renvoyés par compute_pareto / compute_frequency_pareto : 1=catégorie,
    2=valeur, 4=Cumul_%, et l'éventuelle colonne ajoutée par add_seuil80_column.

    C'est ce graphique natif (et non une simple image collée) qui rend le fichier
    Excel exporté fidèle au format attendu par l'entreprise : un vrai Pareto
    (barres + courbe cumulée + seuil 80%) que l'utilisateur peut rouvrir et
    modifier directement dans Excel.
    """
    from openpyxl.chart import BarChart, LineChart, Reference

    last_row = header_row + n_rows

    bar = BarChart()
    bar.type = "col"
    bar.title = title
    bar.style = 10
    bar.y_axis.title = value_axis_title
    bar.y_axis.majorGridlines = None
    bar.gapWidth = 40
    bar.x_axis.delete = False

    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    bar.add_data(val_ref, titles_from_data=True)
    bar.set_categories(cat_ref)

    line = LineChart()
    cumul_ref = Reference(worksheet, min_col=cumul_col, min_row=header_row, max_row=last_row)
    line.add_data(cumul_ref, titles_from_data=True)
    line.y_axis.axId = 200
    line.y_axis.title = "Cumul %"
    line.y_axis.scaling.min = 0
    line.y_axis.scaling.max = 100
    line.y_axis.crosses = "max"
    line.y_axis.majorGridlines = None

    if seuil_col is not None:
        seuil_ref = Reference(worksheet, min_col=seuil_col, min_row=header_row, max_row=last_row)
        line.add_data(seuil_ref, titles_from_data=True)

    cumul_series = line.series[0]
    cumul_series.marker.symbol = "circle"
    cumul_series.graphicalProperties.line.solidFill = "C00000"
    cumul_series.graphicalProperties.line.width = 20000
    cumul_series.smooth = False

    if seuil_col is not None:
        seuil_series = line.series[1]
        seuil_series.marker.symbol = "none"
        seuil_series.graphicalProperties.line.solidFill = "999999"
        seuil_series.graphicalProperties.line.dashStyle = "dash"
        seuil_series.smooth = False

    bar += line
    bar.width = width
    bar.height = height
    worksheet.add_chart(bar, anchor)


def add_repartition_chart_excel(worksheet, header_row, n_rows, cat_col, val_col, title,
                                 anchor, width=11, height=8, donut=True):
    """
    Insère un graphique de répartition (donut par défaut, comme le donut "Répartition
    du downtime par famille" du fichier Cutting Analysis) à partir d'un tableau déjà
    écrit dans la feuille, en-tête à `header_row`, `n_rows` lignes de données.
    """
    from openpyxl.chart import PieChart, DoughnutChart, Reference

    chart = DoughnutChart() if donut else PieChart()
    chart.title = title
    last_row = header_row + n_rows
    val_ref = Reference(worksheet, min_col=val_col, min_row=header_row, max_row=last_row)
    cat_ref = Reference(worksheet, min_col=cat_col, min_row=header_row + 1, max_row=last_row)
    chart.add_data(val_ref, titles_from_data=True)
    chart.set_categories(cat_ref)
    chart.width = width
    chart.height = height
    worksheet.add_chart(chart, anchor)


def summary_kpis(df: pd.DataFrame) -> dict:
    """Indicateurs clés affichés en haut du dashboard."""
    return {
        "nb_evenements": len(df),
        "duree_totale_min": round(df["Duree_min"].sum(), 1),
        "duree_totale_h": round(df["Duree_min"].sum() / 60, 1),
        "machine_top": df.groupby("Machine")["Duree_min"].sum().idxmax() if not df.empty else "N/A",
        "cause_top": df.groupby("Cause")["Duree_min"].sum().idxmax() if not df.empty else "N/A",
        "periode_debut": df["Date"].min(),
        "periode_fin": df["Date"].max(),
    }


# ---------------------------------------------------------------------------
# EXPORT POWERPOINT — reprend TOUS les graphiques (images PNG déjà générées par
# l'app avec Plotly/kaleido, donc mêmes couleurs et mêmes titres) et tableaux du
# dashboard dans une présentation prête à partager, sans avoir à ouvrir Excel.
# ---------------------------------------------------------------------------

THEME_COLOR = "1F4E78"      # même bleu que l'en-tête des feuilles Excel exportées
ACCENT_COLOR = "C0392B"     # même rouge que l'en-tête du Plan d'action Excel
LIGHT_BG = "F2F6FB"


def _pptx_add_table(slide, df, left, top, width, max_rows=12, header_color=THEME_COLOR):
    """Ajoute un tableau PowerPoint natif (pas une image) à partir d'un DataFrame,
    avec le même en-tête coloré/gras que les feuilles Excel exportées."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    df_show = df.head(max_rows)
    n_rows, n_cols = df_show.shape[0] + 1, df_show.shape[1]
    row_h = Inches(0.35)
    height = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    header_rgb = RGBColor.from_string(header_color)
    for j, col_name in enumerate(df_show.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i in range(df_show.shape[0]):
        for j in range(df_show.shape[1]):
            cell = table.cell(i + 1, j)
            val = df_show.iloc[i, j]
            cell.text = "" if pd.isna(val) else str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)

    bottom = top + height
    if len(df) > max_rows:
        note = slide.shapes.add_textbox(left, bottom + Inches(0.03), width, Inches(0.3))
        note.text_frame.text = (
            f"… {len(df) - max_rows} ligne(s) supplémentaire(s) — voir l'export Excel "
            f"pour le détail complet."
        )
        note.text_frame.paragraphs[0].font.size = Pt(9)
        note.text_frame.paragraphs[0].font.italic = True
        note.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        bottom += Inches(0.3)
    return bottom


def build_pptx_report(sections: list, title: str, subtitle: str = None,
                       closing_text: str = None) -> "BytesIO":
    """
    Construit une présentation PowerPoint (.pptx) prête à télécharger, reprenant
    TOUS les graphiques et tableaux déjà calculés/affichés par le dashboard —
    mêmes titres, mêmes légendes, mêmes couleurs (bleu Versigent #1F4E78, rouge
    plan d'action #C0392B) — pour que l'équipe puisse partager les résultats en
    réunion sans avoir à ouvrir Excel.

    `sections` est une liste de dicts décrivant chaque slide de contenu, avec les
    clés optionnelles suivantes :
      - "title"   : titre de la slide (str) — reprend en général le même texte
                    que le st.header()/st.subheader() correspondant dans l'app.
      - "caption" : texte d'explication sous le titre (str) — reprend en général
                    le même texte que le st.caption() correspondant.
      - "kpis"    : liste de tuples (label, valeur) affichés en grandes cartes
                    (utilisé pour la slide de synthèse).
      - "image"   : bytes PNG d'un graphique déjà généré par l'app (fig.to_image).
      - "image2"  : bytes PNG d'un second graphique, affiché à côté du premier
                    (ex : Pareto + camembert de répartition, comme dans l'app).
      - "table"   : DataFrame affiché en tableau natif PowerPoint (tronqué aux
                    premières lignes si besoin, avec une note pour indiquer que
                    le détail complet est dans l'export Excel).
      - "table_header_color" : couleur d'en-tête du tableau pour cette slide
                    (ex : rouge #C0392B pour la slide Plan d'action, comme dans
                    l'export Excel).

    Retourne un BytesIO prêt à passer à st.download_button.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from io import BytesIO

    THEME = RGBColor.from_string(THEME_COLOR)
    ACCENT = RGBColor.from_string(ACCENT_COLOR)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    LIGHT = RGBColor.from_string(LIGHT_BG)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    SLIDE_W = prs.slide_width

    def _bleed_rect(slide, color_rgb):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    # --- Slide de titre (fond bleu foncé, même identité que le dashboard) ---
    slide = prs.slides.add_slide(blank_layout)
    _bleed_rect(slide, THEME)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(38)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.9))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(16)
        tf2.paragraphs[0].font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)

    # --- Slides de contenu ---
    for section in sections:
        slide = prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.65))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.text = section.get("title", "")
        ttf.paragraphs[0].font.size = Pt(24)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = THEME

        y = Inches(1.0)
        if section.get("caption"):
            cap_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(0.8))
            ctf = cap_box.text_frame
            ctf.word_wrap = True
            ctf.text = section["caption"]
            ctf.paragraphs[0].font.size = Pt(12)
            ctf.paragraphs[0].font.italic = True
            ctf.paragraphs[0].font.color.rgb = RGBColor(0x60, 0x60, 0x60)
            y = Inches(1.85)

        if section.get("kpis"):
            kpis_list = section["kpis"]
            n = len(kpis_list)
            gap = Inches(0.25)
            total_w = Inches(12.3)
            box_w = Emu(int((total_w - gap * (n - 1)) / n))
            x = Inches(0.5)
            for label, value in kpis_list:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, Inches(1.5))
                card.fill.solid()
                card.fill.fore_color.rgb = LIGHT
                card.line.color.rgb = THEME
                card.line.width = Pt(1)
                card.shadow.inherit = False
                vtf = card.text_frame
                vtf.word_wrap = True
                vtf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
                vtf.text = str(value)
                vtf.paragraphs[0].font.size = Pt(24)
                vtf.paragraphs[0].font.bold = True
                vtf.paragraphs[0].font.color.rgb = THEME
                vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
                p2 = vtf.add_paragraph()
                p2.text = label
                p2.font.size = Pt(11)
                p2.font.color.rgb = DARK
                p2.alignment = PP_ALIGN.CENTER
                x = Emu(int(x) + int(box_w) + int(gap))
            y = y + Inches(1.75)

        image_bytes = section.get("image")
        image2_bytes = section.get("image2")
        table_df = section.get("table")
        header_color = section.get("table_header_color", THEME_COLOR)

        if image_bytes and image2_bytes:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.4), y, width=Inches(6.1))
            slide.shapes.add_picture(BytesIO(image2_bytes), Inches(6.75), y, width=Inches(6.1))
        elif image_bytes and table_df is not None:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.4), y, width=Inches(7.4))
            _pptx_add_table(slide, table_df, Inches(8.0), y, Inches(4.9), max_rows=10,
                             header_color=header_color)
        elif image_bytes:
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(1.15), y, width=Inches(11.0))
        elif table_df is not None:
            bottom = _pptx_add_table(slide, table_df, Inches(0.7), y, Inches(11.9), max_rows=14,
                                      header_color=header_color)
            y = bottom + Inches(0.15)

        if section.get("body_text"):
            body_box = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.9), Inches(7.3) - y)
            btf = body_box.text_frame
            btf.word_wrap = True
            btf.text = section["body_text"]
            btf.paragraphs[0].font.size = Pt(15)
            btf.paragraphs[0].font.color.rgb = DARK
            btf.paragraphs[0].line_spacing = 1.25

    # --- Slide de clôture (même fond bleu que la slide de titre) ---
    if closing_text:
        slide = prs.slides.add_slide(blank_layout)
        _bleed_rect(slide, THEME)
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(3.1), Inches(11.0), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = closing_text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
